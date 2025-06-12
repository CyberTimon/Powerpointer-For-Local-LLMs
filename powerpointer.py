import collections.abc
from pptx import Presentation
from pptx.util import Inches
import requests
import json
import re
import random
import time
import prompts

# These are all available designs:
#Design 1 = Envelope, beige
#Design 2 = Blue Bubble
#Design 3 = Light Blue Black
#Design 4 = Black, dark
#Design 5 = wood
#Design 6 = Multicolored, Simple
#Design 7 = Black, white

# Language Model Configuration
API_HOST = "127.0.0.1"
API_PORT = 5005
API_URL = f"http://{API_HOST}:{API_PORT}/v1/chat/completions"
ADDITIONAL_LLM_INSTRUCTION = "" # Can be "/no_think" for QWEN thinking models, so they don't output thinking tags (which breaks the pptx generator).

def llm_generate(messages):
    payload = {'messages': messages, 'stream': False}
    headers = {'Content-Type': 'application/json'}
    try:
        response = requests.post(API_URL, json=payload, headers=headers)
        return response.json() if response.status_code == 200 else {'error': True, 'message': f'API error: {response.status_code}'}
    except Exception as e:
        return {'error': True, 'message': f'API error: {str(e)}'}

def create_ppt_text(prompt, slides, info=""):
    final_prompt = prompts.make_prompt(prompt, slides, info)

    messages = [
        {"role": "system", "content": "You are a helpful assistant that creates PowerPoint presentations."},
        {"role": "user", "content": final_prompt}
    ]
    
    response = llm_generate(messages)

    if isinstance(response, dict) and response.get('error', False):
        return "Title: Error generating presentation content: " + response.get('message', 'Unknown error')

    try:
        raw_result = response['choices'][0]['message']['content']
        result = raw_result.replace("<think>\n\n</think>\n\n", "").strip() # Fix for QWEN thinking models
        return result
    except (KeyError, IndexError) as e:
        return "Title: Error parsing API response: " + str(e)

def create_ppt(text_file, design_number, ppt_name):
    prs = Presentation(f"Designs/Design-{design_number}.pptx")
    slide_count = 0
    header = ""
    content = ""
    last_slide_layout_index = -1
    firsttime = True
    with open(text_file, 'r', encoding='utf-8') as f:
        for line_num, line in enumerate(f):
            if line.startswith('Title:'):
                header = line.replace('Title:', '').strip()
                slide = prs.slides.add_slide(prs.slide_layouts[0])
                title = slide.shapes.title
                title.text = header
                body_shape = slide.shapes.placeholders[1]
                continue
            elif line.startswith('Slide:'):
                if slide_count > 0:
                    slide = prs.slides.add_slide(prs.slide_layouts[slide_layout_index])
                    title = slide.shapes.title
                    title.text = header
                    body_shape = slide.shapes.placeholders[slide_placeholder_index]
                    tf = body_shape.text_frame
                    tf.text = content
                content = "" 
                slide_count += 1
                slide_layout_index = last_slide_layout_index
                layout_indices = [1, 7, 8]
                while slide_layout_index == last_slide_layout_index:
                    if firsttime == True:
                        slide_layout_index = 1
                        slide_placeholder_index = 1
                        firsttime = False
                        break
                    slide_layout_index = random.choice(layout_indices)
                    if slide_layout_index == 8:
                        slide_placeholder_index = 2
                    else:
                        slide_placeholder_index = 1
                last_slide_layout_index = slide_layout_index
                continue
            elif line.startswith('Header:'):
                header = line.replace('Header:', '').strip()
                continue
            elif line.startswith('Content:'):
                content = line.replace('Content:', '').strip()
                next_line = f.readline().strip()
                while next_line and not next_line.startswith('#'):
                    content += '\n' + next_line
                    next_line = f.readline().strip()
                continue
                
        if content:
            slide = prs.slides.add_slide(prs.slide_layouts[slide_layout_index])
            title = slide.shapes.title
            title.text = header
            body_shape = slide.shapes.placeholders[slide_placeholder_index]
            tf = body_shape.text_frame
            tf.text = content
            
    prs.save(f'GeneratedPresentations/{ppt_name}.pptx')
    file_path = f"GeneratedPresentations/{ppt_name}.pptx"
    return f"{file_path}"

def generate_ppt(prompt, add_info, slides, theme):
    prompt = re.sub(r'[^\w\s.\-\(\)]', '', prompt)
    if not theme:
        print("No theme selected, using default theme.")
    if theme > 7:
        theme = 1
        print("Invalid theme number, default theme will be applied.")
    elif theme == 0:
        theme = 1
        print("Invalid theme number, default theme will be applied.")

    print("Generating the powerpoint, this could take some time depending on your gpu...\n")
    
    with open(f'Cache/{prompt}.txt', 'w', encoding='utf-8') as f:
        f.write(create_ppt_text(prompt, slides, add_info))

    ppt_path = create_ppt(f'Cache/{prompt}.txt', theme, prompt)
    return str(ppt_path)

def main():
    print("Welcome to the powerpoint generator!")
    topic = input("Topic for the powerpoint: ")
    add_info = input("Consider this in the powerpoint (enter if none): ")
    if not add_info:
        add_info = ""
    slides = input("Number of slides: ")
    theme = int(input("Select theme of the powerpoint (1-7): "))
    start_time = time.time()
    print ("Generated and saved under:", generate_ppt(topic, add_info, slides, theme))
    end_time = time.time()
    print ("Time used for generating:", round((end_time - start_time), 2))
    
main()
